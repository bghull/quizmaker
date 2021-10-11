import csv

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

TITLE = 0
RULES = 1
ROUND_START = 2
QUESTION = 3
ANSWERS_1 = 4
ANSWERS_2 = 5
ANSWERS_AUDIO = 6
BUMPER_1 = 7
BUMPER_2 = 8
BUMPER_3 = 9
BUMPER_4 = 10
CLOSER = 11

def import_sheet(worksheet_values):
    return {"R" + l[1] + l[2][:1] + l[3]: {"Text": l[4], "Notes": l[5]} for l in worksheet_values[1:]}

def import_csv(infile):
    with open(infile, newline="", encoding="utf-8-sig") as csvfile:
        csvdict = csv.DictReader(csvfile)
        return {"R" + d["Round"] + d["Type"][:1] + d["Number"]: {"Text": d["Text"], "Notes": d["Notes"]} for d in csvdict}

def build_round(prs, number, data, audio=None):
    global bumpers
    slide = prs.slides.add_slide(prs.slide_layouts[ROUND_START])
    slide.placeholders[0].text = f"Round {number}"
    slide.placeholders[10].text = data[f"R{number}C"]["Text"]
    slide.placeholders[11].text = data[f"R{number}D"]["Text"]
    notes = slide.notes_slide.notes_text_frame
    if c_notes := data[f"R{number}C"]["Notes"]:
        notes.text = f"Category: {c_notes}"
    if d_notes := data[f"R{number}D"]["Notes"]:
        p = notes.add_paragraph()
        run = p.add_run()
        run.text = f"Description: {d_notes}"
    if audio: # No question slides for audio rounds
        prs.slides.add_slide(prs.slide_layouts[bumpers.pop()])
        slide = prs.slides.add_slide(prs.slide_layouts[ANSWERS_AUDIO])
        slide.placeholders[0].text = f"Round {number} Answers"
        for a in range(1, 8):
            slide.placeholders[a + 9].text = f"A{a}: " + data[f"R{number}A{a}"]["Text"]
    else:
        # Questions
        for q in range(1, 8):
            try:
                data[f"R{number}Q{q}"]["Text"]
            except KeyError:
                data[f"R{number}Q{q}"] = {"Text": "(No question provided)", "Notes": ""}
            slide = prs.slides.add_slide(prs.slide_layouts[QUESTION])
            slide.placeholders[0].text = f"Question {q}"
            slide.placeholders[10].text = data[f"R{number}Q{q}"]["Text"]
            notes = slide.notes_slide.notes_text_frame
            if q_notes := data[f"R{number}Q{q}"]["Notes"]:
                notes.text = f"Q{q}: {q_notes}"
            if a_notes := data[f"R{number}A{q}"]["Notes"]:
                p = notes.add_paragraph()
                run = p.add_run()
                run.text = f"A{q}: {a_notes}"
            #slide.placeholders[10].text_frame.fit_text(max_size=60)
        prs.slides.add_slide(prs.slide_layouts[bumpers.pop()])
        # Answers 1-4
        slide = prs.slides.add_slide(prs.slide_layouts[ANSWERS_1])
        slide.placeholders[0].text = f"Round {number} Answers"
        tf = slide.placeholders[10].text_frame
        notes = slide.notes_slide.notes_text_frame
        for q in range(1, 5):
            if q == 1:
                tf.paragraphs[0].text = f"Q{q}: {data[f'R{number}Q{q}']['Text']}"
                if q_notes := data[f"R{number}Q{q}"]["Notes"]:
                    notes.text = f"Q{q}: {q_notes}"
            else:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"Q{q}: {data[f'R{number}Q{q}']['Text']}"
                if q_notes := data[f"R{number}Q{q}"]["Notes"]:
                    p = notes.add_paragraph()
                    run = p.add_run()
                    run.text = f"Q{q}: {q_notes}"
            p = tf.add_paragraph()
            run = p.add_run()
            if q == 4:
                run.text = f"A{q}: {data[f'R{number}A{q}']['Text']}"
            else:
                run.text = f"A{q}: {data[f'R{number}A{q}']['Text']}\n"
            run.font.bold = True
            run.font.size = Pt(36)
            run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            if a_notes := data[f"R{number}A{q}"]["Notes"]:
                p = notes.add_paragraph()
                run = p.add_run()
                run.text = f"A{q}: {a_notes}"

        # Answers 5-7
        slide = prs.slides.add_slide(prs.slide_layouts[ANSWERS_2])
        slide.placeholders[0].text = f"Round {number} Answers (cont.)"
        tf = slide.placeholders[10].text_frame
        notes = slide.notes_slide.notes_text_frame
        for q in range(5, 8):
            if q == 5:
                tf.paragraphs[0].text = f"Q{q}: {data[f'R{number}Q{q}']['Text']}"
                if q_notes := data[f"R{number}Q{q}"]["Notes"]:
                    notes.text = f"Q{q}: {q_notes}"
            else:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"Q{q}: {data[f'R{number}Q{q}']['Text']}"
                if q_notes := data[f"R{number}Q{q}"]["Notes"]:
                    p = notes.add_paragraph()
                    run = p.add_run()
                    run.text = f"Q{q}: {q_notes}"
            p = tf.add_paragraph()
            run = p.add_run()
            if q == 7:
                run.text = f"A{q}: {data[f'R{number}A{q}']['Text']}"
            else:
                run.text = f"A{q}: {data[f'R{number}A{q}']['Text']}\n"
            run.font.bold = True
            run.font.size = Pt(36)
            run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
            if a_notes := data[f"R{number}A{q}"]["Notes"]:
                p = notes.add_paragraph()
                run = p.add_run()
                run.text = f"A{q}: {a_notes}"

def build_quiz(template, data, audio=None, sheet=None):
    global bumpers
    bumpers = [BUMPER_1, BUMPER_4, BUMPER_3, BUMPER_2, BUMPER_1]
    prs = Presentation(template)
    if sheet:
        all_data = import_sheet(data)
    else:
        all_data = import_csv(data)
    prs.slides.add_slide(prs.slide_layouts[RULES])
    build_round(prs, 1, all_data)
    build_round(prs, 2, all_data)
    if audio:
        build_round(prs, 3, all_data, audio=True)
    else:
        build_round(prs, 3, all_data)
    build_round(prs, 4, all_data)
    build_round(prs, 5, all_data)
    prs.slides.add_slide(prs.slide_layouts[CLOSER])
    return prs
